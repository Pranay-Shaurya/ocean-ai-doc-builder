import io
from typing import Literal

from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt

from database import Project


def export_to_docx(project: Project) -> bytes:
    doc = Document()
    doc.add_heading(project.title, level=0)
    doc.add_paragraph(f"Topic: {project.topic}")

    for section in project.sections:
        doc.add_heading(section.heading, level=1)
        for paragraph in (section.content or "").split("\n"):
            doc.add_paragraph(paragraph.strip())

    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer.read()


def export_to_pptx(project: Project) -> bytes:
    presentation = Presentation()

    # Title slide
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = project.title
    slide.placeholders[1].text = project.topic

    bullet_layout = presentation.slide_layouts[1]
    for section in project.sections:
        slide = presentation.slides.add_slide(bullet_layout)
        slide.shapes.title.text = section.heading
        body = slide.shapes.placeholders[1].text_frame
        body.text = ""
        for line in (section.content or "").split("\n"):
            text = line.strip()
            if not text:
                continue
            if not body.text:
                body.text = text
            else:
                body.add_paragraph().text = text
        for paragraph in body.paragraphs:
            paragraph.font.size = Pt(18)

    buffer = io.BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer.read()


def export_project(project: Project, file_type: Literal["word", "ppt"]) -> tuple[bytes, str]:
    if file_type == "word":
        data = export_to_docx(project)
        mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    else:
        data = export_to_pptx(project)
        mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    return data, mime

